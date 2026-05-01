#!/usr/bin/env python3
"""
Generate a draft PPTX deck for each talk in src/content/talks/, by opening
talks/jmaclabs-template-final.pptx and adding slides via the template's own
slide layouts. The template's master, theme, logo, fonts and layout styling
are preserved exactly because we are working inside the template file.

For each talk:
  1. Start from the template (preserves master + all 18 layouts).
  2. Remove the template's example slides.
  3. For every [SLIDE NN — kind: payload] cue in the talk script, add a slide
     using the layout best suited to the cue's kind, and populate placeholders.
  4. Save as talks/<slug>/deck-draft.pptx.

Cue → layout mapping:

  BLACK / cold open      → Layout  6  Blank
  TITLE_CARD / outro     → Layout  0  Title Slide        (CENTER_TITLE + SUBTITLE)
  SECTION (part header)  → Layout  2  Section Header     (TITLE + BODY)
  PULL_QUOTE             → Layout 11  Quote with Caption (TITLE + quote BODY + attrib BODY)
  STAT (big number)      → Layout 10  Title and Caption  (big TITLE + supporting BODY)
  LIST                   → Layout  1  Title and Content  (TITLE + bulleted OBJECT)
  SPLIT (two-column)     → Layout  4  Comparison         (TITLE + L header/body + R header/body)
  IMAGE                  → Layout  8  Picture with Caption  (TITLE + PICTURE + caption)
  OUTRO                  → Layout 12  Name Card
  GENERIC fallback       → Layout  5  Title Only

The template's layouts already carry the JMacLabs branding: logo placement,
type system (Plus Jakarta Sans / Space Grotesk), surface and accent colors,
slide numbering, and footer chrome. This script does not redraw any of that.
"""

from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

REPO = Path("/Users/jmac/git/personal/jmaclabs")
TEMPLATE = REPO / "talks" / "jmaclabs-template-final.pptx"
TALKS_SRC = REPO / "src" / "content" / "talks"
TALKS_OUT = REPO / "talks"

# Layout index constants (verified by inspecting the template):
LO_TITLE_SLIDE       = 0
LO_TITLE_AND_CONTENT = 1
LO_SECTION_HEADER    = 2
LO_TWO_CONTENT       = 3
LO_COMPARISON        = 4
LO_TITLE_ONLY        = 5
LO_BLANK             = 6
LO_PICTURE_CAPTION   = 8
LO_TITLE_CAPTION     = 10
LO_QUOTE_CAPTION     = 11
LO_NAME_CARD         = 12

CUE_RE = re.compile(
    r"`?\[\s*SLIDE\s+(?P<num>\d+)\s+[—\-]\s+(?P<body>[^`\]]+?)\]`?",
    re.IGNORECASE,
)

# ---------------------------------------------------------------------------
# Cue parsing / classification
# ---------------------------------------------------------------------------

def extract_cues(md: str) -> list[dict]:
    return [{"num": int(m.group("num")), "body": m.group("body").strip()}
            for m in CUE_RE.finditer(md)]


def classify(cue: dict) -> dict:
    body = cue["body"]
    low = body.lower()

    if low.startswith("black slide") or low == "black" or low.startswith("blank"):
        return {"kind": "BLACK"}

    if low.startswith("title card"):
        m = re.search(r'"(.+)"', body)
        text = m.group(1) if m else body.split(":", 1)[-1].strip()
        return {"kind": "TITLE_CARD", "text": text}

    if low.startswith("outro") or low.startswith("closing card") or low.startswith("final card"):
        m = re.search(r'"(.+)"', body)
        text = m.group(1) if m else body.split(":", 1)[-1].strip()
        return {"kind": "OUTRO", "text": text}

    if low.startswith("pull-quote") or low.startswith("pull quote") or low.startswith("bold statement"):
        m = re.search(r'"(.+)"', body)
        text = m.group(1) if m else body.split(":", 1)[-1].strip()
        return {"kind": "PULL_QUOTE", "text": text}

    if low.startswith("title:"):
        m = re.search(r'"(.+)"', body)
        text = m.group(1) if m else body.split(":", 1)[-1].strip()
        return {"kind": "SECTION", "text": text}

    if low.startswith("stat callout") or low.startswith("stat:"):
        m = re.search(r'"(.+)"', body)
        text = m.group(1) if m else body.split(":", 1)[-1].strip()
        return {"kind": "STAT", "text": text}

    # Bio-strip / role bullets
    if low.startswith("short bio") or low.startswith("bio strip") or low.startswith("bio bullets"):
        items_str = body.split(":", 1)[-1]
        items = [s.strip(" \"'") for s in re.split(r"\s*[·•]\s*", items_str) if s.strip()]
        return {"kind": "LIST", "items": items, "title": "John MacDonald"}

    # List builds / checklists / list reveals
    if (low.startswith("list builds") or low.startswith("checklist build")
            or low.startswith("list reveals") or low.startswith("checklist with")
            or low.startswith("list:")):
        items_str = body.split(":", 1)[-1]
        # split on " · " or " 1) " style markers
        items = [s.strip(" \"'") for s in re.split(r"\s*[·•]\s*|\s*\d+\)\s*", items_str) if s.strip()]
        return {"kind": "LIST", "items": items}

    # Two-column comparisons
    if (low.startswith("split") or low.startswith("two columns") or low.startswith("two-cost")
            or "↔" in body or low.startswith("two ")):
        m = re.findall(r'"([^"]+)"', body)
        if len(m) >= 2:
            return {"kind": "SPLIT", "left": m[0], "right": m[1]}
        payload = body.split(":", 1)[-1].strip()
        parts = re.split(r"\s+/\s+|\s+↔\s+", payload)
        if len(parts) >= 2:
            return {"kind": "SPLIT", "left": parts[0].strip(' "'),
                    "right": parts[1].strip(' "')}

    # Timeline as list
    if low.startswith("timeline"):
        items_str = body.split(":", 1)[-1]
        items = [s.strip(" \"'") for s in re.split(r"\s+/\s+|\s*·\s*|\s+;\s+", items_str) if s.strip()]
        return {"kind": "LIST", "items": items, "title": "Timeline"}

    # Series recap → static list
    if low.startswith("recap") or low.startswith("series recap"):
        return {"kind": "LIST", "title": "Series so far", "items": [
            "Part 1: Bubbles Can Build Foundations",
            "Part 2: The Rewiring of Retail",
            "Part 3: From DVD Mail to Streaming",
            "Part 4: The Internet Got Ambient",
            "Part 5: Wireless and the Skipped Stack",
        ]}

    # Image cue
    if (low.startswith("image") or low.startswith("logo") or low.startswith("visual")
            or low.startswith("close-up") or low.startswith("aerial") or low.startswith("split illustration")
            or low.startswith("split image") or low.startswith("diagram")):
        caption = body.split(":", 1)[-1].strip(' "') if ":" in body else body
        return {"kind": "IMAGE", "caption": caption}

    return {"kind": "GENERIC", "text": body}


# ---------------------------------------------------------------------------
# Placeholder helpers
# ---------------------------------------------------------------------------

def find_ph(slide, idx: int):
    """Return placeholder by index, or None."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def set_ph_text(slide, idx: int, text: str):
    ph = find_ph(slide, idx)
    if ph is None or not ph.has_text_frame:
        return False
    tf = ph.text_frame
    tf.text = text
    return True


def set_ph_bullets(slide, idx: int, items: list[str]):
    ph = find_ph(slide, idx)
    if ph is None or not ph.has_text_frame:
        return False
    tf = ph.text_frame
    # First item replaces the existing single paragraph
    tf.text = items[0] if items else ""
    for item in items[1:]:
        p = tf.add_paragraph()
        p.text = item
    return True


def remove_existing_slides(prs: Presentation):
    """Remove slides that exist in the template, leaving an empty deck.
       Layouts and master remain because they live on the slide_master."""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        sldIdLst.remove(sldId)
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass


# ---------------------------------------------------------------------------
# Builders that *use* the template's layouts
# ---------------------------------------------------------------------------

def add_black(prs):
    return prs.slides.add_slide(prs.slide_layouts[LO_BLANK])


def add_title_card(prs, title: str, subtitle: str = ""):
    s = prs.slides.add_slide(prs.slide_layouts[LO_TITLE_SLIDE])
    set_ph_text(s, 0, title)              # CENTER_TITLE
    set_ph_text(s, 1, subtitle or "John MacDonald · jmaclabs.com")  # SUBTITLE
    return s


def add_section(prs, text: str):
    s = prs.slides.add_slide(prs.slide_layouts[LO_SECTION_HEADER])
    set_ph_text(s, 0, text)
    set_ph_text(s, 1, "")
    return s


def add_pull_quote(prs, quote: str, attribution: str = "John MacDonald"):
    s = prs.slides.add_slide(prs.slide_layouts[LO_QUOTE_CAPTION])
    # Layout 11 has TITLE (idx 0), and BODY placeholders at idx 2 and idx 13.
    # We treat idx 2 as the quote and idx 13 as the attribution; if the
    # template authored these the other way around the user can swap in PPT.
    set_ph_text(s, 0, "")  # leave the title field empty for a clean quote
    set_ph_text(s, 2, f"“{quote}”")
    set_ph_text(s, 13, f"— {attribution}")
    return s


def add_stat(prs, text: str):
    """Big-number slide. Split text on the first ' — ', ' : ', or ' / '
       so the first segment becomes the title (huge number) and the rest
       becomes the caption."""
    s = prs.slides.add_slide(prs.slide_layouts[LO_TITLE_CAPTION])
    parts = re.split(r"\s*[—:]\s*", text, maxsplit=1)
    big = parts[0].strip(' "')
    small = parts[1].strip(' "') if len(parts) == 2 else ""
    set_ph_text(s, 0, big)
    set_ph_text(s, 2, small)
    return s


def add_list(prs, items: list[str], title: str = ""):
    s = prs.slides.add_slide(prs.slide_layouts[LO_TITLE_AND_CONTENT])
    set_ph_text(s, 0, title)
    set_ph_bullets(s, 1, items if items else [""])
    return s


def add_split(prs, left: str, right: str, title: str = ""):
    s = prs.slides.add_slide(prs.slide_layouts[LO_COMPARISON])
    # Layout 4 has TITLE (0), BODY headers (1, 3), and OBJECT bodies (2, 4).
    # We use the BODY headers to carry the short label.
    set_ph_text(s, 0, title)
    set_ph_text(s, 1, left)
    set_ph_text(s, 2, "")
    set_ph_text(s, 3, right)
    set_ph_text(s, 4, "")
    return s


def add_image(prs, caption: str, title: str = ""):
    s = prs.slides.add_slide(prs.slide_layouts[LO_PICTURE_CAPTION])
    # Layout 8: TITLE (0), PICTURE (1), BODY caption (2).
    # We leave the PICTURE placeholder empty so PowerPoint shows the
    # "Click icon to add picture" affordance. Caption carries the cue
    # description so the user knows what art to drop in.
    set_ph_text(s, 0, title or "")
    set_ph_text(s, 2, caption)
    return s


def add_outro(prs, text: str):
    s = prs.slides.add_slide(prs.slide_layouts[LO_NAME_CARD])
    set_ph_text(s, 0, "Thanks for watching.")
    set_ph_text(s, 2, text or "jmaclabs.com  ·  Read the full essay at the link in the description")
    return s


def add_generic(prs, text: str):
    s = prs.slides.add_slide(prs.slide_layouts[LO_TITLE_ONLY])
    set_ph_text(s, 0, text)
    return s


# ---------------------------------------------------------------------------
# Per-talk driver
# ---------------------------------------------------------------------------

def build_deck(talk_md_path: Path) -> Path:
    md = talk_md_path.read_text()
    title_match = re.search(r"^title:\s*\"?(.+?)\"?\s*$", md, re.MULTILINE)
    title = title_match.group(1).strip() if title_match else talk_md_path.stem
    cues = extract_cues(md)
    if not cues:
        return None

    slug = re.sub(r"^\d{4}-\d{2}-\d{2}-", "", talk_md_path.stem)
    out_dir = TALKS_OUT / slug
    out_dir.mkdir(exist_ok=True)
    out_path = out_dir / "deck-draft.pptx"

    # Open the template (carries master, theme, layouts, logo placement)
    prs = Presentation(str(TEMPLATE))
    remove_existing_slides(prs)

    for cue in cues:
        info = classify(cue)
        kind = info["kind"]
        try:
            if kind == "BLACK":
                add_black(prs)
            elif kind == "TITLE_CARD":
                add_title_card(prs, info.get("text", title))
            elif kind == "SECTION":
                add_section(prs, info["text"])
            elif kind == "PULL_QUOTE":
                add_pull_quote(prs, info["text"])
            elif kind == "STAT":
                add_stat(prs, info["text"])
            elif kind == "LIST":
                add_list(prs, info["items"], info.get("title", ""))
            elif kind == "SPLIT":
                add_split(prs, info["left"], info["right"])
            elif kind == "IMAGE":
                add_image(prs, info["caption"])
            elif kind == "OUTRO":
                add_outro(prs, info.get("text", ""))
            else:
                add_generic(prs, info.get("text", cue["body"]))
        except Exception as e:
            print(f"    error on cue {cue['num']} ({kind}): {e}")
            add_generic(prs, cue["body"])

    prs.save(out_path)
    return out_path


def main():
    if not TEMPLATE.exists():
        print(f"Template not found at {TEMPLATE}", file=sys.stderr)
        sys.exit(1)
    talks = sorted(TALKS_SRC.glob("*.md"))
    print(f"Template: {TEMPLATE.relative_to(REPO)}")
    print(f"Talks:    {len(talks)}\n")
    for tf in talks:
        print(tf.name)
        out = build_deck(tf)
        if out:
            print(f"  → {out.relative_to(REPO)}")


if __name__ == "__main__":
    main()
